import imaplib
import email
import email.utils
import email.header
import time
import uuid
import smtplib
import PyPDF2
import docx2txt

from rtf_converter import rtf_to_txt
from pptx import Presentation
from datetime import datetime
from htmldocx import HtmlToDocx

import markdown

from email import encoders
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.header import decode_header

import os
import json

import constants
import main
import ocr_pdf

MODE = 'EMAIL'  # EMAIL or MICROPHONE or FORCE AUDIO or FORCE TEXT

FORCED_TEXT_FILENAME = "transcription.txt"
FORCED_AUDIO_FILENAME = "54 Clay Brook Rd 2.m4a"
WORK_TO_DO_DIR = "work-to-do"

email_pass = main.email_pass
email_user = main.email_user


def write_text_file(text, filepath):
    with open(filepath, "w") as f:
        f.write(text)
    return


def convert_pdf_to_txt(pdf_path):
    with open(pdf_path, 'rb') as file:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(file)

        # Initialize text variable to store all content
        plain_text = ""

        # Loop through all pages and extract text
        for i in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[i]
            plain_text += page.extract_text()

    return plain_text


def convert_rtf_to_txt(rtf_path):
    with open(rtf_path, 'r', encoding='utf-8') as file:
        rtf_content = file.read()

    # Convert RTF content to plain text
    plain_text = rtf_to_txt(rtf_content)
    return plain_text


def convert_pptx_to_text(pptx_path):
    prs = Presentation(pptx_path)
    plain_text = ""

    # Iterate through each slide and its shapes to extract text
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                plain_text += shape.text + "\n"

    return plain_text


def write_json_from_text_filepath(from_email, subject, detail_level, text, text_filepath):
    if text_filepath:
        with open(text_filepath, "r") as f:
            text = f.read()

    json_dict = {
        "from": from_email,
        "subject": subject,
        "text": text,
        "detail_level": detail_level
    }

    guid_str = str(uuid.uuid4())
    guid_filename = f"file_{guid_str}.json"
    guid_filepath = WORK_TO_DO_DIR + "/" + guid_filename

    # Writing data to a JSON file using json.dump()
    with open(guid_filepath, "w") as json_file:
        json.dump(json_dict, json_file)

    return guid_filepath


def send_error_response_and_cleanup(filepath, work_filepath, from_email):
    send_email(from_email,
               "Invalid Request",
               "Request must have one and only one M4A, WAV, PDF, RTF, PPTX, DOC or TXT attachment "
               "and a valid subject line.",
               []
               )
    if os.path.exists(filepath):
        os.remove(filepath)

    if os.path.exists(work_filepath):
        os.remove(work_filepath)


def convert_markdown_to_docx(subject, has_evaluation, outline):
    filename_list = []

    with open(main.SUMMARY_FILENAME, "r") as f:
        summary_txt_file_contents = [line for line in f]

    subject_name = main.determine_subject_name(subject, summary_txt_file_contents)

    preprocess(main.SUMMARY_FILENAME, subject_name)
    summary_html_file = convert_markdown_to_html(main.SUMMARY_FILENAME)
    summary_docx_file = convert_html_to_docx(summary_html_file, subject_name) + ".docx"

    filename_list.append(summary_docx_file)

    if has_evaluation:
        preprocess(main.EVALUATION_FILENAME, subject_name)
        evaluation_html_file = convert_markdown_to_html(main.EVALUATION_FILENAME)
        evaluation_docx_file = convert_html_to_docx(evaluation_html_file, subject_name) + ".docx"

        overall_conclusion = get_overall_conclusion(main.EVALUATION_FILENAME)

        filename_list.append(evaluation_docx_file)

    else:
        overall_conclusion = ""

    email_body_txt = "See attachments.\n\n"
    email_body_txt += overall_conclusion + '\n\n'
    email_body_txt += 'Thank you for using Investment Evaluator. To learn more, contact lars@perpetual-labs.com'

    return filename_list, email_body_txt


def preprocess(filename, subject_name):
    with open(filename, "r") as f:
        file_contents = [line for line in f]

    todays_datetime = datetime.now().strftime("%Y-%m-%d %H:%M")
    title = f"**{filename.split('.')[0].upper()} OF {subject_name}    {todays_datetime}**\n\n"

    modified_lines = []
    for line in file_contents:
        p = line.find(":**")
        if p != -1:
            heading = line[:p].strip() + "**\n\n"
            line = line[p + 3:]
            while line.startswith("*"):
                line = line[1:]
            line = heading + line

        modified_lines.append(line)

    with open(filename, "w") as f:
        f.write(title)
        for line in modified_lines:
            f.write(line)

    return


def convert_markdown_to_html(filename):
    filepath = ""
    if os.path.exists(filename):
        with open(filename, "r") as f:
            text = f.read()

        html = markdown.markdown(text)
        filepath = os.path.splitext(filename)[0] + ".html"

        with open(filepath, 'w') as f:
            f.write(html)

    return filepath


def convert_html_to_docx(filename, subject_name):
    output_filename = os.path.splitext(filename)[0].upper() + f'-{subject_name}'
    htmlconverter = HtmlToDocx()
    htmlconverter.parse_html_file(filename, output_filename)

    return output_filename


def get_overall_conclusion(evaluation_file):
    overall_conclusion = ""
    if os.path.exists(evaluation_file):
        with open(evaluation_file, "r") as f:
            text = f.read().replace("#", "").replace("*", "")

        p = text.find("OVERALL CONCLUSION")
        if p != -1:
            overall_conclusion = text[p + 19:].strip()
            if overall_conclusion.startswith(":"):
                overall_conclusion = overall_conclusion[1:].strip()

    return overall_conclusion


def get_emails_and_create_work_files():
    no_new_work_to_do = True
    while no_new_work_to_do:
        # Connect to Gmail's IMAP server
        imap_server = imaplib.IMAP4_SSL("imap.gmail.com")
        imap_server.login(email_user, email_pass)
        imap_server.select('inbox')

        # Search for all unread emails
        status, email_ids = imap_server.search(None, 'UNSEEN')

        if status == 'OK':
            email_id_list = email_ids[0].split()
            unread_count = len(email_id_list)

            if unread_count > 0:
                main.log_message(f"Number of unread emails: {unread_count}")
                attachment_dir = 'email_attachments'
                if not os.path.exists(attachment_dir):
                    os.mkdir(attachment_dir)

                # Fetch and download attachments from each unread email
                work_filepath = ""

                for email_id in email_id_list:
                    valid_attachments = 0
                    _, msg_data = imap_server.fetch(email_id, '(RFC822)')
                    raw_email = msg_data[0][1]
                    msg = email.message_from_bytes(raw_email)
                    from_email = email.utils.parseaddr(msg.get("From"))[1]

                    try:
                        subject, encoding = decode_header(msg["Subject"])[0]

                    except:
                        subject = ""
                        encoding = None

                    if isinstance(subject, bytes):
                        subject = subject.decode(encoding or "utf-8")

                    subject = subject.strip()
                    while subject.find(":") != -1:
                        subject = subject[subject.find(":")+1:].strip()

                    subject = subject.upper() or "DEFAULT"

                    detail_level = 5
                    if subject[:8] == "SUMMARY-":
                        p = subject.find("-")
                        if p != -1:
                            detail_level = subject[p+1:p+3]
                        subject = "SUMMARY"

                    if subject not in constants.summary_prompts.keys():
                        subject = "DEFAULT"

                    if subject in constants.summary_prompts.keys():

                        # Process attachments
                        msg_txt = ""
                        for part in msg.walk():
                            if part.get_content_maintype() == 'multipart':
                                continue
                            if part.get('Content-Disposition') is None:
                                msg_txt = msg_txt or part.get_payload()
                                continue

                            filename = part.get_filename()
                            filename = filename.replace("\r", "").replace("\n", "")

                            if filename:
                                decoded_header = email.header.decode_header(filename)
                                decoded_filename = ''
                                for filename_part, charset in decoded_header:
                                    if charset:
                                        decoded_filename += filename_part.decode(charset)
                                    else:
                                        decoded_filename += filename_part

                                filename = decoded_filename.upper()

                            filepath = ""

                            if filename:
                                if filename.endswith((".M4A", ".WAV", ".TXT", ".PDF", ".RTF", ".PPTX", ".DOCX", ".XL")):
                                    valid_attachments += 1
                                    if valid_attachments == 1:
                                        filepath = os.path.join(attachment_dir, filename)
                                        with open(filepath, 'wb') as f:
                                            f.write(part.get_payload(decode=True))
                                        main.log_message(f"Downloaded attachment: {filename}")

                                        root_filepath, _ = os.path.splitext(filepath)
                                        transcription_filename = root_filepath + ".TXT"

                                        if filename.endswith((".M4A", ".WAV")):
                                            main.log_message("transcribe start")
                                            main.transcribe_audio_using_aai(filepath, transcription_filename)
                                            main.log_message("transcribe end")

                                        elif filename.endswith(".PDF"):
                                            text = ocr_pdf.ocr_pdf(filepath)
                                            write_text_file(text, transcription_filename)

                                        elif filename.endswith(".RTF"):
                                            text = convert_rtf_to_txt(filepath)
                                            write_text_file(text, transcription_filename)

                                        elif filename.endswith(".PPTX"):
                                            text = convert_pptx_to_text(filepath)
                                            write_text_file(text, transcription_filename)

                                        elif filename.endswith(".DOCX"):
                                            text = docx2txt.process(filepath)
                                            write_text_file(text, transcription_filename)

                                        elif filename.endswith(".TXT"):
                                            transcription_filename = filepath

                                        elif filename.endswith(".XL"):
                                            text = ocr_pdf.load_and_convert_from_s3(filepath)
                                            write_text_file(text, transcription_filename)

                                        work_filepath = write_json_from_text_filepath(
                                            from_email,
                                            subject,
                                            detail_level,
                                            "",
                                            transcription_filename
                                        )

                                        if os.path.exists(filepath):
                                            os.remove(filepath)
                                        no_new_work_to_do = False

                                        main.log_message("we got some work to do (attachment) ...")

                                    else:
                                        send_error_response_and_cleanup(
                                            filepath,
                                            work_filepath,
                                            from_email
                                        )
                            else:
                                send_error_response_and_cleanup(filepath, work_filepath, from_email)

                        if len(msg_txt) > 500 and no_new_work_to_do:
                            work_filepath = write_json_from_text_filepath(
                                from_email,
                                "SUMMARY",
                                detail_level,
                                msg_txt,
                                ""
                            )

                            no_new_work_to_do = False
                            main.log_message("we got some work to do (email text) ...")

                    else:
                        send_error_response_and_cleanup("none", work_filepath, from_email)
        else:
            main.log_message("Failed to retrieve unread emails.")

        try:
            imap_server.logout()

        except:
            pass

        if no_new_work_to_do:
            time.sleep(30)

    return


def send_email(recipient_email, subject, body, attachments):
    # Set the sender and recipient email addresses
    sender_email = email_user
    receiver_email = recipient_email
    password = email_pass

    # Set the subject and body of the email
    # subject = "Evaluation of Investment Opportunity"
    # body = "See attachments."

    # Create a multipart email
    outbound_email = MIMEMultipart()
    outbound_email["From"] = sender_email
    outbound_email["To"] = receiver_email
    outbound_email["Subject"] = subject
    outbound_email.attach(MIMEText(body, "plain"))

    for attachment_file in attachments:
        if attachment_file:
            # Open the file in binary mode and create a MIME object
            with open(attachment_file, "rb") as attachment:
                part = MIMEBase("application", "octet-stream")
                part.set_payload(attachment.read())

            # Encode the payload and add the necessary headers
            encoders.encode_base64(part)
            part.add_header(
                "Content-Disposition",
                f"attachment; filename= {attachment_file}",
            )

            # Attach the MIME object to the email
            outbound_email.attach(part)

    # Set up the server and port
    server = smtplib.SMTP('smtp.gmail.com', 587)

    # Log into the server
    server.starttls()  # Start the TLS (Transport Layer Security) mode to encrypt the session
    server.login(sender_email, password)

    # Send the email
    text = outbound_email.as_string()  # Convert the MIMEMultipart object to a string
    server.sendmail(sender_email, receiver_email, text)

    # Quit the server
    server.quit()

